import streamlit as st
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline
import torch
import base64
import os
import io
from pptx import Presentation
import moviepy.editor as mp
from PIL import Image, ImageDraw, ImageFont
import pyttsx3
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
import textwrap
from pptx.util import Pt
from openai_pipeline import openai_pipeline
import openai

st.set_page_config(layout="wide")

# Sidebar: OpenAI API key
st.sidebar.header("Options")
use_openai = st.sidebar.checkbox("Use OpenAI?")
api_key = None
if use_openai:
    api_key = st.sidebar.text_input("Enter your OpenAI API key", type="password")

# -------------------------
# Model and tokenizer loading
# -------------------------
checkpoint = "t5-small"

tokenizer = AutoTokenizer.from_pretrained(checkpoint)
base_model = AutoModelForSeq2SeqLM.from_pretrained(checkpoint, dtype=torch.float32)

# Create a summarization pipeline once
summarizer = pipeline(
    "summarization",
    model=base_model,
    tokenizer=tokenizer,
    device=0 if torch.cuda.is_available() else -1,
    max_length=250,
    min_length=200
)

def file_preprocessing(file_path):
    loader = PyPDFLoader(file_path)
    pages = loader.load_and_split()  
    return pages  # Return list of pages directly

# -------------------------
# Pipeline for summarization
# -------------------------
def llm_pipeline(file_path):
    if api_key:
        try:
            summaries, titles = openai_pipeline(file_path, api_key)
        except:
            st.error("❌ Invalid OpenAI API key!")
            return [], []
    
    else:
        pages = file_preprocessing(file_path)
        summaries = []
        titles = []
        for i, page in enumerate(pages):
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=300, chunk_overlap=80)
            chunks = text_splitter.split_documents([page])
            
            # Merge chunks of this page into one string
            page_text = "".join([chunk.page_content for chunk in chunks])
            
            # Summarize this page
            summary = summarizer(page_text)
            summaries.append(summary[0]['summary_text'])
            titles.append(f"Slide {i+1}")
    
    return summaries, titles  #List of summaries, one per page

# -------------------------
# Display PDF in Streamlit
# -------------------------
@st.cache_data
def displayPDF(file_path):
    with open(file_path, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')
    pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="600" type="application/pdf"></iframe>'
    st.markdown(pdf_display, unsafe_allow_html=True)


def generate_slides(summaries, titles, template_file=None, output_file="final_presentation.pptx"):
    if template_file and os.path.exists(template_file):
        prs = Presentation(template_file)
    else:
        prs = Presentation()

    for i, text in enumerate(summaries):
        # Use the template's first slide layout or pick specific layout index
        slide_layout = prs.slide_layouts[1]  
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        slide.shapes.title.text = titles[i]
        

        # Add body text
        textbox = slide.placeholders[1]
        textbox.text = text

        # Optional: adjust font size for long text
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)  # adjust as needed

    prs.save(output_file)
    return output_file


def generate_slide_images(texts, titles, output_prefix="slide", output_dir="slide_images", font_path="arial.ttf"):
    os.makedirs(output_dir, exist_ok=True)
    image_files = []

    # Fonts
    title_font = ImageFont.truetype(font_path, 60)
    body_font = ImageFont.truetype(font_path, 45)

    for i, text in enumerate(texts):
        # Create blank slide
        img = Image.new("RGB", (1280, 720), color="white")
        draw = ImageDraw.Draw(img)

        # Add slide title
        title = titles[i]
       
        draw.text((50, 50), title, fill="black", font=title_font)

        # Wrap body text
        lines = []
        for paragraph in text.split("\n"):
            wrapped = textwrap.wrap(paragraph, width=40)  # adjust width as needed
            lines.extend(wrapped if wrapped else [""])

        # Draw body text
        y_text = 150
        for line in lines:
            draw.text((50, y_text), line, fill="black", font=body_font)
            bbox = draw.textbbox((0, 0), line, font=body_font)
            line_height = bbox[3] - bbox[1]
            y_text += line_height + 10

        filename = os.path.join(output_dir, f"{output_prefix}_{i}.png")
        img.save(filename)
        image_files.append(filename)

    return image_files


def generate_audio_clips(summaries, titles, output_prefix="audio"):
    engine = pyttsx3.init()
    audio_files = []

    for i, text in enumerate(summaries):
        filename = f"{output_prefix}_{i}.mp3"
        engine.save_to_file(titles[i] + text, filename)
        audio_files.append(filename)
    
    engine.runAndWait()
    return audio_files

def slides_to_video(slide_images, audio_files, output_file="slides_video.mp4"):
    clips = []
    for img, audio in zip(slide_images, audio_files):
        audio_clip = AudioFileClip(audio)
        slide_clip = ImageClip(img).set_duration(audio_clip.duration).set_audio(audio_clip)
        clips.append(slide_clip)
    
    final = concatenate_videoclips(clips, method="compose")
    final.write_videofile(output_file, fps=1)

# -------------------------
# Streamlit App
# -------------------------
st.set_page_config(layout="wide")

def main():
    st.title("PDF Documents to Summarised Slides")

    uploaded_file = st.file_uploader("Upload your PDF file", type=['pdf'])

    if uploaded_file is not None:
        # Ensure data folder exists and save uploaded file
        os.makedirs("data", exist_ok=True)
        filepath = os.path.join("data", uploaded_file.name)
        with open(filepath, "wb") as temp_file:
            temp_file.write(uploaded_file.read())

        col1, col2 = st.columns(2)
        with col1:
            st.info("Uploaded File")
            displayPDF(filepath)

        # Initialize session_state only once
        if "summary_texts" not in st.session_state:
            st.session_state.summary_texts = []  # will hold editable summaries
        if "summarized" not in st.session_state:
            st.session_state.summarized = False  # flag to avoid re-running summarizer

        with col2:
            # Only run summarization once
            if st.button("Summarize") or st.session_state.summarized:
                if not st.session_state.summarized:
                    summaries, titles = llm_pipeline(filepath)
                    st.session_state.summary_texts = summaries
                    st.session_state.titles = titles
                    st.session_state.summarized = True


                # Display editable text areas
                for i in range(len(st.session_state.summary_texts)):
                    st.session_state.summary_texts[i] = st.text_area(
                        label=f"Page {i+1} Summary",
                        value=st.session_state.summary_texts[i],
                        height=150,
                        key=f"page_summary_{i}"
                    )

            # Buttons side by side for Video and Slides
            btn_col1, btn_col2 = st.columns(2)

            with btn_col1:
                if st.button("Generate Video"):
                    # 1. Generate slide images directly from summaries
                    if not st.session_state.summary_texts:
                        st.error("Error!")
                    else:
                        slide_images = generate_slide_images(st.session_state.summary_texts, st.session_state.titles)

                        # 2. Generate audio clips for each slide
                        audio_files = generate_audio_clips(st.session_state.summary_texts, st.session_state.titles)

                        # 3. Combine into video
                        output_video = "generated_slides_video.mp4"
                        slides_to_video(slide_images, audio_files, output_file=output_video)

                        # 4. Read video bytes for download
                        with open(output_video, "rb") as f:
                            video_bytes = f.read()

                        st.download_button(
                            label="Download Video",
                            data=video_bytes,
                            file_name="generated_slides_video.mp4",
                            mime="video/mp4"
                    )

            with btn_col2:
                if st.button("Generate Slides"):
                    if not st.session_state.summary_texts:
                        st.error("Error!")
                    else:
                        pptx_file = generate_slides(st.session_state.summary_texts, st.session_state.titles)
                        with open(pptx_file, "rb") as f:
                            st.download_button(
                                label="Download Slides",
                                data=f,
                                file_name=pptx_file,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )

if __name__ == "__main__":
    main()



